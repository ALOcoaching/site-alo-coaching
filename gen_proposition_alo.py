#!/usr/bin/env python3
"""
Génère la proposition commerciale ALO Coaching V2 — Design amélioré.
Part du contenu de la V1 modifiée par l'utilisateur.
Branding : Navy (#1a2332) + Gold (#c8993e)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

# ── Couleurs ALO ──
NAVY = RGBColor(0x1A, 0x23, 0x32)
NAVY_LIGHT = RGBColor(0x24, 0x2F, 0x40)
GOLD = RGBColor(0xC8, 0x99, 0x3E)
GOLD_LIGHT = RGBColor(0xD4, 0xAF, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF1, 0xEB)  # Warm light beige (more elegant than gray)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MEDIUM_TEXT = RGBColor(0x55, 0x55, 0x55)
LIGHT_TEXT = RGBColor(0x88, 0x88, 0x88)
SUBTLE_BORDER = RGBColor(0xD8, 0xD0, 0xC4)  # Warm border
RED_ACCENT = RGBColor(0xCC, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BASE_DIR = os.path.dirname(__file__)
LOGO_PATH = os.path.join(BASE_DIR, "site-alo-coaching", "img", "logo-alo.png")
LOGO_NOIR_PATH = os.path.join(BASE_DIR, "site-alo-coaching", "img", "logo-alo-noir.png")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ──
def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    # Set corner radius
    sp = shape._element
    sp_pr = sp.find(qn('a:prstGeom'))
    if sp_pr is not None:
        av_lst = sp_pr.find(qn('a:avLst'))
        if av_lst is None:
            av_lst = sp_pr.makeelement(qn('a:avLst'), {})
            sp_pr.append(av_lst)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=14, color=DARK_TEXT, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(4),
                  space_after=Pt(2), level=0, italic=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p

def add_gold_accent(slide, left, top, height, width=Inches(0.06)):
    """Thin elegant gold accent bar"""
    return add_shape(slide, left, top, width, height, fill_color=GOLD)

def add_header_bar(slide, title, subtitle=None):
    """Elegant header with gradient-style navy bar + gold underline"""
    add_shape(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=NAVY)
    # Gold thin line under header
    add_shape(slide, 0, Inches(1.15), SLIDE_W, Pt(3), fill_color=GOLD)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(9), Inches(0.65),
                 title, font_size=28, color=WHITE, bold=True, font_name="Calibri")
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(9), Inches(0.35),
                     subtitle, font_size=13, color=GOLD, font_name="Calibri")

def add_footer(slide, text="ALO Coaching & Formation  |  Arnaud Loyet  |  Proposition commerciale – Mars 2026"):
    add_shape(slide, Inches(0.5), SLIDE_H - Inches(0.5), SLIDE_W - Inches(1), Pt(1.5), fill_color=GOLD)
    add_text_box(slide, Inches(0.5), SLIDE_H - Inches(0.42), SLIDE_W - Inches(1), Inches(0.35),
                 text, font_size=8, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)

def add_logo_noir(slide, left=None, top=Inches(0.22), height=Inches(0.65)):
    """Use black logo on light backgrounds, regular on dark"""
    path = LOGO_NOIR_PATH if os.path.exists(LOGO_NOIR_PATH) else LOGO_PATH
    if not os.path.exists(path):
        return
    if left is None:
        left = SLIDE_W - Inches(2)
    slide.shapes.add_picture(path, left, top, height=height)

def add_logo_blanc(slide, left=None, top=Inches(0.4), height=Inches(0.85)):
    """Regular logo on dark backgrounds"""
    if not os.path.exists(LOGO_PATH):
        return
    if left is None:
        left = SLIDE_W - Inches(2.6)
    slide.shapes.add_picture(LOGO_PATH, left, top, height=height)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — PAGE DE GARDE (design premium)
# ═══════════════════════════════════════════════════════════════
sl1 = prs.slides.add_slide(prs.slide_layouts[6])

# Full navy background
add_shape(sl1, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)

# Subtle darker panel left side for visual depth
add_shape(sl1, 0, 0, Inches(0.6), SLIDE_H, fill_color=RGBColor(0x12, 0x18, 0x24))

# Gold vertical accent
add_shape(sl1, Inches(0.6), Inches(1.8), Inches(0.08), Inches(3.5), fill_color=GOLD)

# Horizontal gold line top
add_shape(sl1, Inches(1.2), Inches(1.5), Inches(8), Pt(1), fill_color=GOLD)

# Logo top-right
add_logo_blanc(sl1, left=SLIDE_W - Inches(3), top=Inches(0.5), height=Inches(1))

# Title
add_text_box(sl1, Inches(1.2), Inches(1.9), Inches(9), Inches(1),
             "PROPOSITION COMMERCIALE", font_size=36, color=GOLD,
             bold=True, font_name="Calibri")

# Subtitle
add_text_box(sl1, Inches(1.2), Inches(2.85), Inches(9), Inches(0.9),
             "Direction de Projets & Management du Site de Lannion",
             font_size=22, color=WHITE, font_name="Calibri")

# Horizontal gold line under subtitle
add_shape(sl1, Inches(1.2), Inches(3.7), Inches(4), Pt(1), fill_color=GOLD)

# Client + Référence
tb_client = add_text_box(sl1, Inches(1.2), Inches(4.1), Inches(9), Inches(0.7),
             "Client : OPEN", font_size=18, color=RGBColor(0xAA, 0xAA, 0xAA),
             font_name="Calibri")
add_paragraph(tb_client.text_frame, "Référence ALOCF_2026001",
              font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA), space_before=Pt(4))

# Author + date + URL
tb_author = add_text_box(sl1, Inches(1.2), Inches(4.9), Inches(9), Inches(0.8),
             "Arnaud Loyet  —  Mars 2026", font_size=15, color=RGBColor(0x77, 0x77, 0x77),
             font_name="Calibri")
add_paragraph(tb_author.text_frame, "https://aloformationcoaching.netlify.app/",
              font_size=11, color=GOLD, italic=True, space_before=Pt(8))

# Gold bottom bar (elegant thin)
add_shape(sl1, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill_color=GOLD)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ═══════════════════════════════════════════════════════════════
sl_som = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(sl_som, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_BG)

add_header_bar(sl_som, "SOMMAIRE", "Proposition commerciale – Mars 2026")
add_logo_noir(sl_som)

# Sommaire items — visual table of contents
som_items = [
    ("01", "Périmètre juin", "4 semaines · 2 jours/semaine (lundis & mardis) · Basé site de Lannion"),
    ("02", "Périmètre juillet / Août", "Interventions ponctuelles de clôture"),
    ("03", "Périmètre septembre → décembre", "2 options de poursuite · 4 mois · Basé site de Lannion"),
    ("04", "Volet financier", "TJM · Portage Jump · Synthèse budgétaire"),
    ("05", "Conditions & Modalités", "Validité · Facturation · Délais · Confidentialité"),
    ("", "Annexes", "Option de portage ITG (alternative)"),
]

som_top = Inches(1.6)
som_left = Inches(1.5)
som_w = Inches(10.5)
item_h = Inches(0.78)

for i, (num, title, desc) in enumerate(som_items):
    y = som_top + i * (item_h + Inches(0.08))
    bg = CARD_BG if i % 2 == 0 else LIGHT_BG
    is_annexe = (num == "")

    # Row card
    add_rounded_shape(sl_som, som_left, y, som_w, item_h,
                      fill_color=bg, line_color=SUBTLE_BORDER, line_width=Pt(0.5))

    # Gold left accent
    add_gold_accent(sl_som, som_left, y, item_h)

    if not is_annexe:
        # Number circle
        add_rounded_shape(sl_som, som_left + Inches(0.2), y + Inches(0.14),
                          Inches(0.5), Inches(0.5), fill_color=NAVY)
        add_text_box(sl_som, som_left + Inches(0.2), y + Inches(0.14),
                     Inches(0.5), Inches(0.5), num,
                     font_size=16, color=GOLD, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        # Annexe marker
        add_rounded_shape(sl_som, som_left + Inches(0.2), y + Inches(0.14),
                          Inches(0.5), Inches(0.5), fill_color=GOLD)
        add_text_box(sl_som, som_left + Inches(0.2), y + Inches(0.14),
                     Inches(0.5), Inches(0.5), "A",
                     font_size=16, color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Title
    add_text_box(sl_som, som_left + Inches(0.9), y + Inches(0.1),
                 Inches(5), Inches(0.4), title,
                 font_size=16, color=NAVY, bold=True)

    # Description
    add_text_box(sl_som, som_left + Inches(0.9), y + Inches(0.42),
                 Inches(8), Inches(0.3), desc,
                 font_size=11, color=MEDIUM_TEXT)

add_footer(sl_som)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — PÉRIMÈTRE JUIN
# ═══════════════════════════════════════════════════════════════
sl2 = prs.slides.add_slide(prs.slide_layouts[6])

# Light warm background
add_shape(sl2, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_BG)

add_header_bar(sl2, "01  ·  PÉRIMÈTRE — JUIN",
               "4 semaines  ·  2 jours / semaine (lundis & mardis)  ·  8 jours  ·  Basé site de Lannion")
add_logo_noir(sl2)

# ── Left column: Site Lannion ──
col1_left = Inches(0.7)
col1_top = Inches(1.55)
col1_w = Inches(5.6)

# Card with subtle shadow effect (darker shape behind)
add_rounded_shape(sl2, col1_left + Inches(0.03), col1_top + Inches(0.03),
                  col1_w, Inches(4.55), fill_color=RGBColor(0xE0, 0xDA, 0xD0))
# Main card
card1 = add_rounded_shape(sl2, col1_left, col1_top, col1_w, Inches(4.55),
                          fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))

# Card header
add_shape(sl2, col1_left, col1_top, col1_w, Inches(0.6), fill_color=NAVY)
add_text_box(sl2, col1_left + Inches(0.25), col1_top + Inches(0.05), col1_w - Inches(0.5), Inches(0.5),
             "🏢  SITE LANNION", font_size=18, color=GOLD, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

# Gold accent
add_gold_accent(sl2, col1_left, col1_top + Inches(0.6), Inches(3.95))

tb1 = add_text_box(sl2, col1_left + Inches(0.25), col1_top + Inches(0.8),
                   col1_w - Inches(0.5), Inches(3.5), "", font_size=15)
tf1 = tb1.text_frame
tf1.paragraphs[0].text = ""

items_lannion = [
    ("Gestion CDP & Hebdo Prod", "Pilotage de la réunion hebdomadaire de production,\nsuivi des plans d'action, reporting managérial"),
    ("Gestion du site", "Administration générale du site (~1 jour/mois réparti),\ncoordination des équipes Lannion"),
]
for i, (title, desc) in enumerate(items_lannion):
    add_paragraph(tf1, f"▸  {title}", font_size=15, color=NAVY, bold=True,
                  space_before=Pt(16 if i > 0 else 6))
    for line in desc.split("\n"):
        add_paragraph(tf1, f"    {line}", font_size=12, color=MEDIUM_TEXT,
                      space_before=Pt(3))

add_paragraph(tf1, "", font_size=8)
# Volume badge
add_rounded_shape(sl2, col1_left + Inches(0.2), col1_top + Inches(3.65), Inches(2.8), Inches(0.5),
                  fill_color=NAVY)
add_text_box(sl2, col1_left + Inches(0.2), col1_top + Inches(3.65), Inches(2.8), Inches(0.5),
             "≈ 0.5 jour / semaine", font_size=13, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ── Right column: Orange ──
col2_left = Inches(6.8)
col2_w = Inches(5.8)

# Shadow
add_rounded_shape(sl2, col2_left + Inches(0.03), col1_top + Inches(0.03),
                  col2_w, Inches(4.55), fill_color=RGBColor(0xE0, 0xDA, 0xD0))
# Card
card2 = add_rounded_shape(sl2, col2_left, col1_top, col2_w, Inches(4.55),
                          fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))

# Card header
add_shape(sl2, col2_left, col1_top, col2_w, Inches(0.6), fill_color=NAVY)
add_text_box(sl2, col2_left + Inches(0.25), col1_top + Inches(0.05), col2_w - Inches(0.5), Inches(0.5),
             "📡  ORANGE", font_size=18, color=GOLD, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

add_gold_accent(sl2, col2_left, col1_top + Inches(0.6), Inches(3.95))

tb2 = add_text_box(sl2, col2_left + Inches(0.25), col1_top + Inches(0.8),
                   col2_w - Inches(0.5), Inches(3.5), "", font_size=15)
tf2 = tb2.text_frame
tf2.paragraphs[0].text = ""

items_orange = [
    ("delivery  projets Orange", "Suivi des projets MOBA, PDC et CTOOP (planification, staffing, hebdo)\nAccompagnement et coaching des CP du programme Orange du CDP"),
    ("Suivi Delivery Commerce", "Coordination delivery des projets Orange,\nsuivi des engagements et livrables"),
    ("Gestion COPIL clés", "Préparation et animation des comités de pilotage\nstratégiques du périmètre"),
    ("Synchros CDP / DP / Exec", "Points récurrents avec les chefs de projets,\ndirecteurs de projets et direction"),
]
for i, (title, desc) in enumerate(items_orange):
    add_paragraph(tf2, f"▸  {title}", font_size=15, color=NAVY, bold=True,
                  space_before=Pt(12 if i > 0 else 6))
    for line in desc.split("\n"):
        add_paragraph(tf2, f"    {line}", font_size=12, color=MEDIUM_TEXT,
                      space_before=Pt(3))

# Volume badge
add_rounded_shape(sl2, col2_left + Inches(0.2), col1_top + Inches(3.65), Inches(2.8), Inches(0.5),
                  fill_color=NAVY)
add_text_box(sl2, col2_left + Inches(0.2), col1_top + Inches(3.65), Inches(2.8), Inches(0.5),
             "≈ 1.5 jours / semaine", font_size=13, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom summary bar
add_rounded_shape(sl2, Inches(0.7), Inches(6.4), SLIDE_W - Inches(1.4), Inches(0.55),
                  fill_color=NAVY)
add_text_box(sl2, Inches(0.7), Inches(6.4), SLIDE_W - Inches(1.4), Inches(0.55),
             "TOTAL JUIN  ·  8 jours sur 4 semaines  ·  2 jours / semaine (lundis & mardis)",
             font_size=15, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

add_footer(sl2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — JUILLET / AOÛT (clôtures)
# ═══════════════════════════════════════════════════════════════
sl3 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(sl3, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_BG)

add_header_bar(sl3, "02  ·  PÉRIMÈTRE — JUILLET / AOÛT",
               "Interventions ponctuelles de clôture")
add_logo_noir(sl3)

# Central card — wider and more elegant
card_w = Inches(9)
card_left = (SLIDE_W - card_w) / 2
card_top = Inches(1.8)

# Shadow
add_rounded_shape(sl3, card_left + Inches(0.04), card_top + Inches(0.04),
                  card_w, Inches(3.3), fill_color=RGBColor(0xE0, 0xDA, 0xD0))
# Card
add_rounded_shape(sl3, card_left, card_top, card_w, Inches(3.3),
                  fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))
add_gold_accent(sl3, card_left, card_top, Inches(3.3))

tb3 = add_text_box(sl3, card_left + Inches(0.4), card_top + Inches(0.3),
                   card_w - Inches(0.8), Inches(2.8), "", font_size=15)
tf3 = tb3.text_frame
tf3.paragraphs[0].text = ""

add_paragraph(tf3, "📅  Clôture fin juillet", font_size=19, color=NAVY, bold=True,
              space_before=Pt(4))
add_paragraph(tf3, "     Intervention de 2 à 3 jours en fin de mois", font_size=14,
              color=MEDIUM_TEXT, space_before=Pt(6))
add_paragraph(tf3, "     Revue de production, bilan mensuel, préparation de la période estivale",
              font_size=12, color=MEDIUM_TEXT, space_before=Pt(2))

add_paragraph(tf3, "", font_size=10, space_before=Pt(14))

add_paragraph(tf3, "📅  Clôture fin août", font_size=19, color=NAVY, bold=True,
              space_before=Pt(8))
add_paragraph(tf3, "     Intervention de 2 à 3 jours en fin de mois", font_size=14,
              color=MEDIUM_TEXT, space_before=Pt(6))
add_paragraph(tf3, "     Bilan estival, reprise du suivi avec les CP & les équipes,",
              font_size=12, color=MEDIUM_TEXT, space_before=Pt(2))
add_paragraph(tf3, "     préparation de la rentrée, alignement des plannings",
              font_size=12, color=MEDIUM_TEXT, space_before=Pt(2))

# Summary boxes — 3 elegant KPI boxes
box_w = Inches(2.7)
box_h = Inches(1.0)
box_top = card_top + Inches(3.6)
box_gap = Inches(0.3)
total_boxes_w = 3 * box_w + 2 * box_gap
boxes_left = (SLIDE_W - total_boxes_w) / 2

for i, (label, value) in enumerate([
    ("VOLUME MINIMUM", "4 jours\n(2j × 2 mois)"),
    ("VOLUME MAXIMUM", "6 jours\n(3j × 2 mois)"),
    ("PÉRIODE", "Dernière\nsemaine du mois")
]):
    bx = boxes_left + i * (box_w + box_gap)
    add_rounded_shape(sl3, bx, box_top, box_w, box_h, fill_color=NAVY)
    add_text_box(sl3, bx, box_top + Inches(0.08), box_w, Inches(0.25),
                 label, font_size=10, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl3, bx, box_top + Inches(0.35), box_w, Inches(0.6),
                 value, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

add_footer(sl3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — SEPTEMBRE → DÉCEMBRE (2 options)
# ═══════════════════════════════════════════════════════════════
sl4 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(sl4, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_BG)

add_header_bar(sl4, "03  ·  PÉRIMÈTRE — SEPTEMBRE À DÉCEMBRE",
               "2 options de poursuite  ·  4 mois  ·  Basé site de Lannion")
add_logo_noir(sl4)

# ── Option 1 ──
opt1_left = Inches(0.7)
opt1_top = Inches(1.5)
opt1_w = Inches(5.6)

# Shadow + card
add_rounded_shape(sl4, opt1_left + Inches(0.03), opt1_top + Inches(0.03),
                  opt1_w, Inches(5), fill_color=RGBColor(0xE0, 0xDA, 0xD0))
add_rounded_shape(sl4, opt1_left, opt1_top, opt1_w, Inches(5),
                  fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))

# Header bar
add_shape(sl4, opt1_left, opt1_top, opt1_w, Inches(0.55), fill_color=NAVY)
add_text_box(sl4, opt1_left, opt1_top, opt1_w, Inches(0.55),
             "OPTION 1  —  Site / DCDP + DP Orange", font_size=16, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_gold_accent(sl4, opt1_left, opt1_top + Inches(0.55), Inches(4.45))

tb4a = add_text_box(sl4, opt1_left + Inches(0.25), opt1_top + Inches(0.75),
                    opt1_w - Inches(0.5), Inches(3.8), "", font_size=14)
tf4a = tb4a.text_frame
tf4a.paragraphs[0].text = ""

add_paragraph(tf4a, "2 jours / semaine (lundis & mardis)", font_size=17, color=NAVY, bold=True, space_before=Pt(4))
add_paragraph(tf4a, "", font_size=6)

add_paragraph(tf4a, "🏢  Site Lannion", font_size=15, color=NAVY, bold=True, space_before=Pt(8))
add_paragraph(tf4a, "  ▸ Gestion CDP & Hebdo Prod", font_size=12, color=MEDIUM_TEXT, space_before=Pt(4))
add_paragraph(tf4a, "  ▸ Gestion du site (~1j/mois)", font_size=12, color=MEDIUM_TEXT, space_before=Pt(2))

add_paragraph(tf4a, "📡  Orange", font_size=15, color=NAVY, bold=True, space_before=Pt(12))
add_paragraph(tf4a, "  ▸ PSH projets (MOBA + PDC + CTOOP)", font_size=12, color=MEDIUM_TEXT, space_before=Pt(4))
add_paragraph(tf4a, "  ▸ Suivi Delivery Commerce", font_size=12, color=MEDIUM_TEXT, space_before=Pt(2))
add_paragraph(tf4a, "  ▸ COPIL clés du périmètre", font_size=12, color=MEDIUM_TEXT, space_before=Pt(2))
add_paragraph(tf4a, "  ▸ Synchros CDP / DP / Exec", font_size=12, color=MEDIUM_TEXT, space_before=Pt(2))

# Volume badge
add_rounded_shape(sl4, opt1_left + Inches(0.2), opt1_top + Inches(4.15), Inches(4.5), Inches(0.45),
                  fill_color=NAVY)
add_text_box(sl4, opt1_left + Inches(0.2), opt1_top + Inches(4.15), Inches(4.5), Inches(0.45),
             "≈ 8 jours / mois  ·  32 jours sur 4 mois", font_size=13, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ── Option 2 ──
opt2_left = Inches(6.8)
opt2_w = Inches(5.8)

add_rounded_shape(sl4, opt2_left + Inches(0.03), opt1_top + Inches(0.03),
                  opt2_w, Inches(5), fill_color=RGBColor(0xE0, 0xDA, 0xD0))
add_rounded_shape(sl4, opt2_left, opt1_top, opt2_w, Inches(5),
                  fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))

# Header bar — gold to differentiate
add_shape(sl4, opt2_left, opt1_top, opt2_w, Inches(0.55), fill_color=GOLD)
add_text_box(sl4, opt2_left, opt1_top, opt2_w, Inches(0.55),
             "OPTION 2  —  Focus DP Orange", font_size=16, color=NAVY, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_gold_accent(sl4, opt2_left, opt1_top + Inches(0.55), Inches(4.45))

tb4b = add_text_box(sl4, opt2_left + Inches(0.25), opt1_top + Inches(0.75),
                    opt2_w - Inches(0.5), Inches(3.8), "", font_size=14)
tf4b = tb4b.text_frame
tf4b.paragraphs[0].text = ""

add_paragraph(tf4b, "Focus Orange uniquement", font_size=17, color=NAVY, bold=True, space_before=Pt(4))
add_paragraph(tf4b, "(hors périmètre site Lannion)", font_size=12, color=MEDIUM_TEXT, space_before=Pt(2))
add_paragraph(tf4b, "", font_size=6)

add_paragraph(tf4b, "📌  Jours fixes", font_size=15, color=NAVY, bold=True, space_before=Pt(8))
add_paragraph(tf4b, "  ▸ 1 jour / semaine (les lundis)", font_size=13, color=MEDIUM_TEXT, space_before=Pt(4))
add_paragraph(tf4b, "     PSH, Delivery Commerce, synchros récurrentes", font_size=11, color=MEDIUM_TEXT, space_before=Pt(2))

add_paragraph(tf4b, "📌  Jours variables", font_size=15, color=NAVY, bold=True, space_before=Pt(14))
add_paragraph(tf4b, "  ▸ ~2 jours / mois", font_size=13, color=MEDIUM_TEXT, space_before=Pt(4))
add_paragraph(tf4b, "     Réunions hors gestion directe, ne pouvant", font_size=11, color=MEDIUM_TEXT, space_before=Pt(2))
add_paragraph(tf4b, "     être planifiées les lundis/mardis :   CPOM, clôture mensuelle, COPIL imposés client", font_size=11, color=MEDIUM_TEXT, space_before=Pt(2))

# Volume badge
add_rounded_shape(sl4, opt2_left + Inches(0.2), opt1_top + Inches(4.15), Inches(4.5), Inches(0.45),
                  fill_color=NAVY)
add_text_box(sl4, opt2_left + Inches(0.2), opt1_top + Inches(4.15), Inches(4.5), Inches(0.45),
             "≈ 6 jours / mois  ·  24 jours sur 4 mois", font_size=13, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Comparison bar
add_rounded_shape(sl4, Inches(0.7), Inches(6.75), (SLIDE_W - Inches(1.4)) / 2 - Inches(0.1), Inches(0.45),
                  fill_color=NAVY)
add_text_box(sl4, Inches(0.7), Inches(6.75), (SLIDE_W - Inches(1.4)) / 2 - Inches(0.1), Inches(0.45),
             "Option 1 : 32 jours (sept→déc)", font_size=14, color=WHITE,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rounded_shape(sl4, Inches(0.7) + (SLIDE_W - Inches(1.4)) / 2 + Inches(0.1), Inches(6.75),
                  (SLIDE_W - Inches(1.4)) / 2 - Inches(0.1), Inches(0.45),
                  fill_color=GOLD)
add_text_box(sl4, Inches(0.7) + (SLIDE_W - Inches(1.4)) / 2 + Inches(0.1), Inches(6.75),
             (SLIDE_W - Inches(1.4)) / 2 - Inches(0.1), Inches(0.45),
             "Option 2 : 24 jours (sept→déc)", font_size=14, color=NAVY, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(sl4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — VOLET FINANCIER (portage Jump — option préférée)
# ═══════════════════════════════════════════════════════════════
sl5 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(sl5, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_BG)

add_header_bar(sl5, "04  ·  VOLET FINANCIER",
               "TJM 800 € HT  ·  Portage via Jump")
add_logo_noir(sl5)

# ── TJM + Jump side by side ──
tjm_left = Inches(0.7)
tjm_top = Inches(1.45)

# TJM box
tjm_w = Inches(4.2)
tjm_h = Inches(1.2)
add_rounded_shape(sl5, tjm_left, tjm_top, tjm_w, tjm_h, fill_color=NAVY)
add_text_box(sl5, tjm_left, tjm_top + Inches(0.08), tjm_w, Inches(0.3),
             "TAUX JOURNALIER MOYEN", font_size=12, color=GOLD, alignment=PP_ALIGN.CENTER)
add_text_box(sl5, tjm_left, tjm_top + Inches(0.4), tjm_w, Inches(0.65),
             "800 € HT / jour", font_size=30, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Jump portage box
jump_left = Inches(5.3)
jump_w = Inches(7.4)
# Shadow
add_rounded_shape(sl5, jump_left + Inches(0.03), tjm_top + Inches(0.03),
                  jump_w, tjm_h, fill_color=RGBColor(0xE0, 0xDA, 0xD0))
add_rounded_shape(sl5, jump_left, tjm_top, jump_w, tjm_h,
                  fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))
add_gold_accent(sl5, jump_left, tjm_top, tjm_h)

tb_jump = add_text_box(sl5, jump_left + Inches(0.25), tjm_top + Inches(0.1),
                       jump_w - Inches(0.5), tjm_h - Inches(0.2), "", font_size=14)
tf_jump = tb_jump.text_frame
tf_jump.paragraphs[0].text = ""
add_paragraph(tf_jump, "PORTAGE VIA JUMP", font_size=15, color=NAVY, bold=True,
              space_before=Pt(0))
add_paragraph(tf_jump, "Surcoût :  + 150 € / mois  (forfait fixe, indépendant du volume)", font_size=13,
              color=GOLD, bold=True, space_before=Pt(6))

# ── Tableau synthèse avec Jump ──
table_top = Inches(3.0)
table_left = Inches(0.7)

cols = [
    ("Période", Inches(3.2)),
    ("Jours", Inches(1.3)),
    ("Montant HT", Inches(2.2)),
    ("+ Jump (150€/mois)", Inches(2.2)),
    ("Total avec Jump", Inches(2.5)),
]

row_h = Inches(0.42)
col_positions = []
cx = table_left
for label, cw in cols:
    col_positions.append((cx, cw))
    cx += cw

data = [
    ("Juin (4 semaines, 2j/sem)", "8 j", "6 400 €", "+ 150 €", "6 550 €"),
    ("Juillet + Août (clôtures)", "4 à 6 j", "3 200 – 4 800 €", "+ 300 €", "3 500 – 5 100 €"),
    ("Sept→Déc · Option 1 (2j/sem)", "32 j", "25 600 €", "+ 600 €", "26 200 €"),
    ("Sept→Déc · Option 2 (focus Orange)", "24 j", "19 200 €", "+ 600 €", "19 800 €"),
]

# Header
for (cx, cw), (label, _) in zip(col_positions, cols):
    add_shape(sl5, cx, table_top, cw, row_h, fill_color=NAVY)
    add_text_box(sl5, cx, table_top, cw, row_h,
                 label, font_size=11, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Rows with alternating colors
for ri, row in enumerate(data):
    ry = table_top + row_h * (ri + 1)
    bg = CARD_BG if ri % 2 == 0 else LIGHT_BG
    for ci, ((cx, cw), val) in enumerate(zip(col_positions, row)):
        if ci == 4:  # Total column highlighted
            fill = NAVY
            txt_color = GOLD
            txt_bold = True
        elif ci == 3:  # Jump surcoût
            fill = bg
            txt_color = GOLD
            txt_bold = True
        else:
            fill = bg
            txt_color = DARK_TEXT
            txt_bold = False
        add_shape(sl5, cx, ry, cw, row_h, fill_color=fill,
                  line_color=SUBTLE_BORDER, line_width=Pt(0.5))
        al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        add_text_box(sl5, cx + Inches(0.1), ry, cw - Inches(0.2), row_h,
                     val, font_size=11, color=txt_color, bold=txt_bold,
                     alignment=al, anchor=MSO_ANCHOR.MIDDLE)

# ── Totaux annuels ──
total_top = table_top + row_h * 5 + Inches(0.2)
total_w = Inches(3.8)
total_h = Inches(1.0)

for i, (label, val) in enumerate([
    ("TOTAL ANNUEL · Opt 1 + Jump", "Juin + clôtures + 2j/sem\n35 750 – 37 350 €"),
    ("TOTAL ANNUEL · Opt 2 + Jump", "Juin + clôtures + focus Orange\n29 350 – 30 950 €"),
]):
    bx = Inches(0.7) + i * (total_w + Inches(0.4))
    add_rounded_shape(sl5, bx, total_top, total_w, total_h, fill_color=NAVY)
    add_text_box(sl5, bx, total_top + Inches(0.06), total_w, Inches(0.25),
                 label, font_size=10, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl5, bx, total_top + Inches(0.35), total_w, Inches(0.6),
                 val, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

add_footer(sl5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — CONDITIONS & MODALITÉS
# ═══════════════════════════════════════════════════════════════
sl_cond = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(sl_cond, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_BG)

add_header_bar(sl_cond, "05  ·  CONDITIONS & MODALITÉS",
               "Validité · Facturation · Délais · Confidentialité")
add_logo_noir(sl_cond)

# ── Left column: Conditions principales ──
cond_left = Inches(0.7)
cond_top = Inches(1.5)
cond_w = Inches(5.8)
cond_h = Inches(5.0)

# Shadow + card
add_rounded_shape(sl_cond, cond_left + Inches(0.03), cond_top + Inches(0.03),
                  cond_w, cond_h, fill_color=RGBColor(0xE0, 0xDA, 0xD0))
add_rounded_shape(sl_cond, cond_left, cond_top, cond_w, cond_h,
                  fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))
add_gold_accent(sl_cond, cond_left, cond_top, cond_h)

tb_cond = add_text_box(sl_cond, cond_left + Inches(0.3), cond_top + Inches(0.2),
                       cond_w - Inches(0.6), cond_h - Inches(0.4), "", font_size=14)
tf_cond = tb_cond.text_frame
tf_cond.paragraphs[0].text = ""

cond_items = [
    ("📋  Validité de la proposition",
     "La présente proposition est valable 30 jours   à compter de sa date d'émission (mars 2026)."),
    ("💰  Facturation",
     "Facturation mensuelle à terme échu, sur la base des jours effectivement réalisés."),
    ("🔔  Délai de prévenance",
     "1 mois de préavis pour toute modification  ou résiliation de la mission."),
    ("🔒  Confidentialité",
     "Les informations échangées dans le cadre de cette mission restent strictement confidentielles."),
    ("🚗  Déplacements",
     "Mission basée sur le site de Lannion.\nEn cas d'intervention sur le site de Rennes (tout ou partie)   application des frais de déplacement selon barème OPEN."),
]

for i, (title, desc) in enumerate(cond_items):
    add_paragraph(tf_cond, title, font_size=15, color=NAVY, bold=True,
                  space_before=Pt(16 if i > 0 else 4))
    for line in desc.split("\n"):
        add_paragraph(tf_cond, f"    {line}", font_size=12, color=MEDIUM_TEXT,
                      space_before=Pt(3))

# ── Right column: Portage & infos pratiques ──
right_left = Inches(6.9)
right_w = Inches(5.7)

# --- Portage box ---
port_top = cond_top
port_h = Inches(2.3)

add_rounded_shape(sl_cond, right_left + Inches(0.03), port_top + Inches(0.03),
                  right_w, port_h, fill_color=RGBColor(0xE0, 0xDA, 0xD0))
add_rounded_shape(sl_cond, right_left, port_top, right_w, port_h,
                  fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))

add_shape(sl_cond, right_left, port_top, right_w, Inches(0.5), fill_color=NAVY)
add_text_box(sl_cond, right_left, port_top, right_w, Inches(0.5),
             "PORTAGE SALARIAL", font_size=14, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

tb_port = add_text_box(sl_cond, right_left + Inches(0.25), port_top + Inches(0.6),
                       right_w - Inches(0.5), port_h - Inches(0.8), "", font_size=12)
tf_port = tb_port.text_frame
tf_port.paragraphs[0].text = ""
add_paragraph(tf_port, "Option recommandée : Jump", font_size=14, color=NAVY, bold=True,
              space_before=Pt(2))
add_paragraph(tf_port, "Forfait fixe : + 150 € / mois", font_size=13, color=GOLD, bold=True,
              space_before=Pt(6))
add_paragraph(tf_port, "", font_size=4)
add_paragraph(tf_port, "Alternative : ITG (+5% du CA)", font_size=12, color=MEDIUM_TEXT,
              space_before=Pt(4))
add_paragraph(tf_port, "Détail en annexe", font_size=11, color=MEDIUM_TEXT, italic=True,
              space_before=Pt(2))

# --- Infos pratiques ---
info_top = port_top + port_h + Inches(0.2)
info_h = Inches(2.5)

add_rounded_shape(sl_cond, right_left + Inches(0.03), info_top + Inches(0.03),
                  right_w, info_h, fill_color=RGBColor(0xE0, 0xDA, 0xD0))
add_rounded_shape(sl_cond, right_left, info_top, right_w, info_h,
                  fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))

add_shape(sl_cond, right_left, info_top, right_w, Inches(0.5), fill_color=NAVY)
add_text_box(sl_cond, right_left, info_top, right_w, Inches(0.5),
             "INFORMATIONS PRATIQUES", font_size=14, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

tb_info = add_text_box(sl_cond, right_left + Inches(0.25), info_top + Inches(0.6),
                       right_w - Inches(0.5), info_h - Inches(0.8), "", font_size=12)
tf_info = tb_info.text_frame
tf_info.paragraphs[0].text = ""
add_paragraph(tf_info, "Consultant : Arnaud Loyet", font_size=14, color=NAVY, bold=True,
              space_before=Pt(2))
add_paragraph(tf_info, "", font_size=4)
add_paragraph(tf_info, "▸  Disponibilité : à partir de juin 2026", font_size=12, color=MEDIUM_TEXT,
              space_before=Pt(4))
add_paragraph(tf_info, "▸  Localisation : Lannion (sur site)", font_size=12, color=MEDIUM_TEXT,
              space_before=Pt(2))
add_paragraph(tf_info, "▸  Déplacements : Lannion (Rennes = frais OPEN)", font_size=12, color=MEDIUM_TEXT,
              space_before=Pt(2))
add_paragraph(tf_info, "▸  Contact : arnaud.loyet@yahoo.fr", font_size=12, color=GOLD,
              space_before=Pt(6))
add_paragraph(tf_info, "▸  Web : aloformationcoaching.netlify.app", font_size=12, color=GOLD,
              space_before=Pt(2))

add_footer(sl_cond)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — PAGE ANNEXES (intercalaire)
# ═══════════════════════════════════════════════════════════════
sl6 = prs.slides.add_slide(prs.slide_layouts[6])

# Full navy background
add_shape(sl6, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)
add_shape(sl6, 0, 0, Inches(0.6), SLIDE_H, fill_color=RGBColor(0x12, 0x18, 0x24))
add_shape(sl6, Inches(0.6), Inches(2.6), Inches(0.08), Inches(2.5), fill_color=GOLD)

# Logo
add_logo_blanc(sl6, left=SLIDE_W - Inches(3), top=Inches(0.5), height=Inches(1))

# Title
add_text_box(sl6, Inches(1.2), Inches(2.9), Inches(9), Inches(1),
             "ANNEXES", font_size=40, color=GOLD, bold=True)

# Author + URL
tb_a = add_text_box(sl6, Inches(1.2), Inches(4.3), Inches(9), Inches(0.8),
                    "Arnaud Loyet  —  Mars 2026", font_size=15, color=RGBColor(0x77, 0x77, 0x77))
add_paragraph(tb_a.text_frame, "https://aloformationcoaching.netlify.app/",
              font_size=11, color=GOLD, italic=True, space_before=Pt(8))

add_shape(sl6, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill_color=GOLD)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — ANNEXE : Option de portage ITG
# ═══════════════════════════════════════════════════════════════
sl7 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(sl7, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_BG)

# Annexe header — slightly different style (lighter header)
add_shape(sl7, 0, 0, SLIDE_W, Inches(1.15), fill_color=RGBColor(0x2A, 0x33, 0x42))
add_shape(sl7, 0, Inches(1.15), SLIDE_W, Pt(3), fill_color=GOLD)
add_text_box(sl7, Inches(0.8), Inches(0.2), Inches(9), Inches(0.65),
             "ANNEXE — OPTION DE PORTAGE ITG", font_size=26, color=RGBColor(0xBB, 0xBB, 0xBB), bold=True)
add_text_box(sl7, Inches(0.8), Inches(0.7), Inches(9), Inches(0.35),
             "Alternative au portage Jump  ·  Surcoût proportionnel de 5 %", font_size=13, color=GOLD)
add_logo_noir(sl7)

# ── Explication ITG ──
itg_left = Inches(0.7)
itg_top = Inches(1.5)
itg_w = Inches(5.5)
itg_h = Inches(2.0)

add_rounded_shape(sl7, itg_left + Inches(0.03), itg_top + Inches(0.03),
                  itg_w, itg_h, fill_color=RGBColor(0xE0, 0xDA, 0xD0))
add_rounded_shape(sl7, itg_left, itg_top, itg_w, itg_h,
                  fill_color=CARD_BG, line_color=SUBTLE_BORDER, line_width=Pt(0.75))
add_gold_accent(sl7, itg_left, itg_top, itg_h)

tb7 = add_text_box(sl7, itg_left + Inches(0.25), itg_top + Inches(0.15),
                   itg_w - Inches(0.5), itg_h - Inches(0.3), "", font_size=14)
tf7 = tb7.text_frame
tf7.paragraphs[0].text = ""
add_paragraph(tf7, "PORTAGE VIA ITG", font_size=18, color=NAVY, bold=True, space_before=Pt(0))
add_paragraph(tf7, "Surcoût :  + 5 % du montant total facturé", font_size=15, color=GOLD, bold=True,
              space_before=Pt(8))
add_paragraph(tf7, "", font_size=4)
add_paragraph(tf7, "▸  Pourcentage proportionnel au chiffre d'affaires", font_size=12, color=MEDIUM_TEXT,
              space_before=Pt(4))

# ── Tableau ITG ──
table7_top = Inches(3.8)
table7_left = Inches(0.7)

cols7 = [
    ("Période", Inches(3.2)),
    ("Jours", Inches(1.1)),
    ("Montant HT", Inches(2.1)),
    ("Surcoût ITG (5%)", Inches(2.1)),
    ("Total avec ITG", Inches(2.5)),
]

row7_h = Inches(0.4)
col7_pos = []
cx7 = table7_left
for label, cw in cols7:
    col7_pos.append((cx7, cw))
    cx7 += cw

data7 = [
    ("Juin (4 semaines, 2j/sem)", "8 j", "6 400 €", "+ 320 €", "6 720 €"),
    ("Juillet + Août (clôtures)", "4 à 6 j", "3 200 – 4 800 €", "+ 160 – 240 €", "3 360 – 5 040 €"),
    ("Sept→Déc · Option 1 (2j/sem)", "32 j", "25 600 €", "+ 1 280 €", "26 880 €"),
    ("Sept→Déc · Option 2 (focus Orange)", "24 j", "19 200 €", "+ 960 €", "20 160 €"),
]

for (cx, cw), (label, _) in zip(col7_pos, cols7):
    add_shape(sl7, cx, table7_top, cw, row7_h, fill_color=NAVY)
    add_text_box(sl7, cx, table7_top, cw, row7_h,
                 label, font_size=10, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

for ri, row in enumerate(data7):
    ry = table7_top + row7_h * (ri + 1)
    bg = CARD_BG if ri % 2 == 0 else LIGHT_BG
    for ci, ((cx, cw), val) in enumerate(zip(col7_pos, row)):
        if ci == 4:
            fill, txt_color, txt_bold = NAVY, GOLD, True
        elif ci == 3:
            fill, txt_color, txt_bold = bg, RED_ACCENT, True
        else:
            fill, txt_color, txt_bold = bg, DARK_TEXT, False
        add_shape(sl7, cx, ry, cw, row7_h, fill_color=fill,
                  line_color=SUBTLE_BORDER, line_width=Pt(0.5))
        al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        add_text_box(sl7, cx + Inches(0.1), ry, cw - Inches(0.2), row7_h,
                     val, font_size=10, color=txt_color, bold=txt_bold,
                     alignment=al, anchor=MSO_ANCHOR.MIDDLE)

# ── Totaux ITG ──
tot7_top = table7_top + row7_h * 5 + Inches(0.15)
tot7_w = Inches(3.5)
tot7_h = Inches(0.85)

for i, (label, val) in enumerate([
    ("TOTAL ANNUEL · Opt 1 + ITG", "36 960 – 38 640 €"),
    ("TOTAL ANNUEL · Opt 2 + ITG", "30 240 – 31 920 €"),
]):
    bx = Inches(0.7) + i * (tot7_w + Inches(0.35))
    add_rounded_shape(sl7, bx, tot7_top, tot7_w, tot7_h, fill_color=NAVY)
    add_text_box(sl7, bx, tot7_top + Inches(0.05), tot7_w, Inches(0.22),
                 label, font_size=9, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl7, bx, tot7_top + Inches(0.3), tot7_w, Inches(0.5),
                 val, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

add_footer(sl7, text="ALO Coaching & Formation  |  Arnaud Loyet  |  Annexe – Mars 2026")


# ═══════════════════════════════════════════════════════════════
# Sauvegarde
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(BASE_DIR, "Proposition_ALO_Coaching_OPEN_2026 V1.00.pptx")
prs.save(output_path)
print(f"[OK] Fichier généré : {output_path}")
